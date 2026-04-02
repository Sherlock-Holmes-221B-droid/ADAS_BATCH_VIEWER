# HtmltoReport.py
# ------------------------------------------------------------
# ADAS Batch Viewer (Streamlit; can be wrapped by pywebview)
# - Drag/drop HTML -> parses scenarios via parser.py
# - Auto-loads Signals.json from same folder (optional)
# - Global filters (speed/overlap + pass/fail/error)
# - Settings (⚙️): tile columns, tile size, selected-info font size
#   + NEW: choose ANY number of tile parameters (multi-line tiles)
# - Donut + Scenario list (colored dots; respects global filters)
# - Clickable tile panel (minimal face; informative hover with units)
# - "Selected" header + 3-column details (font-size controlled)
# - "All Scenarios" table (AgGrid with fallback; anchored)
# - Export PPT: donut + panel + Top Failures + All Passed/Failed (paginated)
# - Footer: Help (reads Batch_readme.md) + About (popover)
# - Compatible with parser returning df  OR (df, cfg)
# ------------------------------------------------------------

import os
import io
import math
import json
import tempfile
from typing import Optional, Tuple, List, Dict

import numpy as np
import pandas as pd
import streamlit as st
import plotly.graph_objects as go

# Optional grid; gracefully fall back if missing or if component errors out
USE_AGGRID = True
try:
    from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode
except Exception:
    USE_AGGRID = False

# Click capture
try:
    from streamlit_plotly_events import plotly_events
except Exception:
    plotly_events = None  # fallback handled below

# Parser
from parser import parse_report

# PPT export
from pptx import Presentation
from pptx.util import Inches, Pt


# -------------------- UI CONFIG --------------------
st.set_page_config(page_title="ADAS Batch Viewer", layout="wide")
st.title("ADAS Batch – Interactive Dashboard")

COLOR_PASS = "#34A853"  # green
COLOR_FAIL = "#EA4335"  # red
COLOR_MIX  = "#F1DC0B"  # amber
COLOR_UNK  = "#9E9E9E"  # gray (used for 'Error' or unknown)


# -------------------- HELPERS --------------------
def color_for_status(s: str) -> str:
    if s == "Pass": return COLOR_PASS
    if s == "Fail": return COLOR_FAIL
    if s == "Mixed": return COLOR_MIX
    if s == "Error": return COLOR_UNK
    return COLOR_UNK

def scenario_state_from_status(status: str) -> str:
    """Map parser status to tri-state for list/filter: Pass/Fail/Error."""
    return status if status in ("Pass", "Fail") else "Error"

def load_signals_json() -> Dict[str, any]:
    """
    Load Signals.json from the same folder as this script.
    Expected keys:
      - "tile_parameters": [list of column names to show on tile]
      - "units": { column_name: " unit" }
    """
    try:
        here = os.path.dirname(os.path.abspath(__file__))
    except NameError:
        here = os.getcwd()
    path = os.path.join(here, "Signals.json")
    if os.path.exists(path):
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)
            # normalize
            tp = data.get("tile_parameters", [])
            if isinstance(tp, str):
                tp = [tp]
            units = data.get("units", {})
            if not isinstance(units, dict):
                units = {}
            return {"tile_parameters": tp, "units": units}
        except Exception:
            pass
    return {"tile_parameters": [], "units": {}}

def make_donut(passed: int, failed: int):
    """Status distribution donut."""
    values = [passed, failed]
    labels = ["Pass", "Fail"]
    colors = [COLOR_PASS, COLOR_FAIL]
    fig = go.Figure(data=[go.Pie(labels=labels, values=values, hole=0.55,
                                 marker=dict(colors=colors),
                                 sort=False, textinfo='percent+label')])
    fig.update_layout(height=320, showlegend=False, margin=dict(l=10, r=10, t=10, b=10),
                      paper_bgcolor="white")
    return fig

def mini_dot_html(status: str) -> str:
    """Small colored dot for scenario list."""
    col = color_for_status(status)
    return f'<span style="display:inline-block;width:10px;height:10px;border-radius:50%;background:{col};margin-right:8px;vertical-align:middle"></span>'

def render_scenario_list(df: pd.DataFrame, max_rows: int = 1500):
    """
    Render a simple list of scenarios with colored dots (respects global filters via df arg).
    """
    if df.empty:
        st.info("No scenarios for current filters.")
        return

    # Sort by state -> speed -> overlap -> name
    order_cols = [c for c in ["state", "ego_speed", "overlap", "scenario_name"] if c in df.columns]
    d = df.sort_values(order_cols).head(max_rows)
    items = []
    for r in d.itertuples():
        dot = mini_dot_html(getattr(r, "state", getattr(r, "status", "Error")))
        items.append(
            f'<div style="padding:2px 0">{dot}<span title="{r.scenario_name}">{r.scenario_name}</span></div>'
        )

    html = f"""
    <div style="max-height:400px; overflow:auto; border:1px solid #eee; border-radius:6px; padding:8px; background:#fff">
      {''.join(items)}
    </div>
    """
    st.markdown(html, unsafe_allow_html=True)

def _fmt_value(v: any) -> str:
    if pd.isna(v): return "—"
    try:
        return f"{float(v):g}"
    except Exception:
        return str(v)

def make_tile_panel(
    df: pd.DataFrame,
    tile_fields: List[str],
    units: Dict[str, str],
    cols: int = 10,
    marker_size: int = 82,
    label_size: int = 14
):
    """
    Build a clean, label-less clickable tile panel (one tile per scenario).
    - Color: state/status (Pass/Fail/Error)
    - Text: first field bold on line1, then each chosen field on its own line
    - Hover: scenario name + "name: value unit" for each field
    Returns (fig, ordered_df) or None if empty
    """
    if df.empty:
        return None

    if not tile_fields:
        # fall back to a minimal default if nothing selected
        tile_fields = [c for c in ["ego_speed", "overlap"] if c in df.columns] or [df.columns[0]]

    d = df.copy()

    # Ensure 'state'
    if "state" not in d.columns and "status" in d.columns:
        d["state"] = d["status"].apply(scenario_state_from_status)

    # Coerce numeric for chosen fields where possible
    for f in tile_fields:
        if f in d.columns:
            d[f] = pd.to_numeric(d[f], errors="coerce")

    # Require at least the first field (primary) to be present
    primary = tile_fields[0]
    if primary not in d.columns:
        return None
    d = d[pd.notna(d[primary])].copy()
    if d.empty:
        return None

    # Stable layout
    sort_cols = [c for c in tile_fields if c in d.columns] + (["scenario_name"] if "scenario_name" in d.columns else [])
    d = d.sort_values(sort_cols)

    # Grid coords
    n = len(d)
    cols = max(1, int(cols))
    rows = math.ceil(n / cols)
    xs, ys = [], []
    for i in range(n):
        col = (i % cols) + 1
        row = rows - (i // cols)   # top-to-bottom
        xs.append(col)
        ys.append(row)

    # Colors
    colors = [color_for_status(s) for s in d["state"]]

    # Build face text: first field bold, others as separate lines
    def face_lines(row) -> List[str]:
        lines = []
        for j, f in enumerate(tile_fields):
            val = _fmt_value(row.get(f))
            if j == 0:
                lines.append(f"<b>{val}</b>")
            else:
                lines.append(val)
        return lines

    labels = ["<br>".join(face_lines(r._asdict() if hasattr(r, "_asdict") else r))
              for _, r in d[tile_fields].iterrows()]

    # Hover lines: "name: value unit"
    def hover_lines(row) -> List[str]:
        hx = [f"<b>{row.get('scenario_name','')}</b>"]
        for f in tile_fields:
            v = _fmt_value(row.get(f))
            unit = units.get(f, "")
            hx.append(f"{f}: {v}{unit}")
        return hx

    hover_text = ["<br>".join(hover_lines(r)) for _, r in d[["scenario_name"] + tile_fields if "scenario_name" in d.columns else tile_fields].iterrows()]

    # customdata is not strictly needed since we preformat hover; we still keep scenario name for click mapping clarity
    customdata = np.stack([d["scenario_name"]], axis=1) if "scenario_name" in d.columns else None

    fig = go.Figure(data=[go.Scatter(
        x=xs, y=ys, mode="markers+text",
        text=labels, textposition="middle center",
        textfont=dict(color="#FFFFFF", size=label_size),
        marker=dict(
            symbol="square",
            size=marker_size,
            color=colors,
            line=dict(color="rgba(0,0,0,0.18)", width=2)
        ),
        hovertemplate="%{text}<extra></extra>",  # we feed prebuilt HTML in 'text' via separate trace? Instead attach hover in text and use hovertext
    )])

    # Use hovertext instead of hovertemplate to display our lines; we keep text for face
    fig.data[0].hovertext = hover_text
    fig.data[0].hovertemplate = "%{hovertext}<extra></extra>"
    if customdata is not None:
        fig.data[0].customdata = customdata

    # Clean canvas
    fig.update_xaxes(visible=False, range=[0.5, cols + 0.5], fixedrange=True)
    fig.update_yaxes(visible=False, range=[0.5, rows + 0.5], fixedrange=True,
                     scaleanchor="x", scaleratio=1)
    # Estimate height
    fig.update_layout(
        height=max(360, math.ceil(rows * (marker_size * 0.95))),
        margin=dict(l=16, r=16, t=10, b=10),
        plot_bgcolor="white",
        paper_bgcolor="white"
    )

    d = d.reset_index(drop=True)
    return fig, d

# -------------------- PPT EXPORT HELPERS --------------------
def _add_table_slide_paginated(prs: Presentation, title: str, df: pd.DataFrame, columns, rows_per_slide: int = 18):
    """
    Add one or more slides with a compact table that fits on the slide.
    Controls column widths and font size to avoid overflow.
    """
    df = df.loc[:, [c for c in columns if c in df.columns]].copy()
    total_rows = len(df)
    if total_rows == 0:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = title
        s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(9.0), Inches(1.0)).text_frame.text = "No rows"
        return

    start = 0
    while start < total_rows:
        chunk = df.iloc[start:start + rows_per_slide]
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = title + ("" if total_rows <= rows_per_slide else f" (rows {start+1}-{min(start+rows_per_slide,total_rows)}/{total_rows})")
        rows = len(chunk) + 1
        cols = len(columns)

        # Table area
        left = Inches(0.4)
        top = Inches(1.2)
        width = Inches(9.6)
        height = Inches(5.0)

        tbl = s.shapes.add_table(rows=rows, cols=cols, left=left, top=top, width=width, height=height).table
        # Column widths (proportional; tune per your taste)
        width_map = {
            "scenario_name": Inches(4.0),
            "status": Inches(1.0),
            "ego_speed": Inches(1.1),
            "overlap": Inches(1.1),
            "impact_speed": Inches(1.3),
            "t_aeb": Inches(1.1)
        }
        for j, col in enumerate(columns):
            try:
                tbl.columns[j].width = width_map.get(col, Inches(1.2))
            except Exception:
                pass

        # Header
        for j, h in enumerate(columns):
            cell = tbl.cell(0, j)
            cell.text = h.replace("_", " ").title()
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True

        # Data
        for i, r in enumerate(chunk.itertuples(index=False), start=1):
            for j, col in enumerate(columns):
                val = getattr(r, col) if hasattr(r, col) else ""
                cell = tbl.cell(i, j)
                cell.text = "" if (val is None or (isinstance(val, float) and pd.isna(val))) else str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)

        start += rows_per_slide

def export_ppt(view_df: pd.DataFrame, pie_fig, panel_fig, file_bytes: io.BytesIO):
    """
    Build PPT:
    - Title
    - Executive Summary (KPIs + donut)
    - Tile panel image
    - Top Failures table
    - All Passed (paginated) + All Failed (paginated)
    """
    # Save chart images using kaleido
    with tempfile.TemporaryDirectory() as td:
        pie_png   = os.path.join(td, "status.png")
        panel_png = os.path.join(td, "panel.png")
        if pie_fig is not None:
            pie_fig.write_image(pie_png, scale=2)
        if panel_fig is not None:
            panel_fig.write_image(panel_png, scale=2)

        prs = Presentation()

        # Title slide
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = "ADAS Batch – Summary"
        s.placeholders[1].text = "Auto-generated from batch HTML"

        # KPI slide
        df = view_df.copy()
        if "state" not in df.columns and "status" in df.columns:
            df["state"] = df["status"].apply(scenario_state_from_status)

        total = len(df)
        passed = int(df["state"].eq("Pass").sum())
        failed = int(df["state"].eq("Fail").sum())

        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.title.text = "Executive Summary"
        tf = s2.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(4.7), Inches(3.2)).text_frame
        if total == 0:
            tf.text = "No data"
        else:
            tf.text = f"Total runs: {total}\nPass: {passed}\nFail: {failed}\nPass rate: {passed/total:.0%}"
        if pie_fig is not None:
            s2.shapes.add_picture(pie_png, Inches(5.6), Inches(1.0), width=Inches(4.5))

        # Tile panel
        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        s3.shapes.title.text = "Scenario Tiles"
        if panel_fig is not None:
            s3.shapes.add_picture(panel_png, Inches(0.6), Inches(1.2), width=Inches(9.2))

        # Top Failures (compact)
        worst = (
            df[df["state"].eq("Fail")]
            .sort_values(["impact_speed"], ascending=False)
            .head(18)[["scenario_name", "ego_speed", "overlap", "impact_speed", "t_aeb"]]
        )
        _add_table_slide_paginated(prs, "Top Failures", worst,
                                   ["scenario_name", "ego_speed", "overlap", "impact_speed", "t_aeb"],
                                   rows_per_slide=18)

        # All Passed / All Failed (paginated)
        all_pass = df[df["state"].eq("Pass")][["scenario_name","status","ego_speed","overlap","impact_speed","t_aeb"]]
        all_fail = df[df["state"].eq("Fail")][["scenario_name","status","ego_speed","overlap","impact_speed","t_aeb"]]

        _add_table_slide_paginated(prs, f"All Passed ({len(all_pass)})", all_pass,
                                   ["scenario_name","status","ego_speed","overlap","impact_speed","t_aeb"],
                                   rows_per_slide=20)
        _add_table_slide_paginated(prs, f"All Failed ({len(all_fail)})", all_fail,
                                   ["scenario_name","status","ego_speed","overlap","impact_speed","t_aeb"],
                                   rows_per_slide=20)

        prs.save(file_bytes)


# -------------------- APP BODY --------------------
# 1) Auto-load Signals.json (same folder)
signals_cfg = load_signals_json()   # {"tile_parameters": [...], "units": {...}}

# 2) Optional: (legacy) dashboard_config uploader still supported if you want to try alt config
# cfg_dashboard = None
# cfg_file = st.file_uploader("Optional (legacy): Load dashboard_config.json", type=["json"], key="cfguploader")
# if cfg_file is not None:
#     try:
#         cfg_dashboard = json.load(cfg_file)
#     except Exception as ex:
#         st.warning(f"Could not parse JSON config: {ex}")

uploaded = st.file_uploader("Drop your batch HTML here", type=["html", "htm"])

if uploaded:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".html") as tmp:
        tmp.write(uploaded.getvalue())
        html_path = tmp.name

    # Parse report (supports both old and new parser signatures)
    # res = None
    # try:
    #     # If your parser supports config, pass the dashboard config
    #     res = parse_report(html_path, config=cfg_dashboard)  # type: ignore[arg-type]
    # except TypeError:
    #     # Old parser signature without 'config'
    #     res = parse_report(html_path)  # type: ignore[call-arg]

    res = parse_report(html_path)

    # Unpack result (df) or (df, cfg_from_parser)
    if isinstance(res, tuple):
        df, cfg_from_parser = res
        # We don't force cfg_dashboard; parser already handled merging if it supported it.
    else:
        df = res

    # Derive 'state' if missing
    if "state" not in df.columns and "status" in df.columns:
        df["state"] = df["status"].apply(scenario_state_from_status)

    # ---- Filters (global) ----
    with st.expander("Filters", expanded=True):
        c1, c2, c3 = st.columns(3)
        with c1:
            # Speed (primary) filter if present
            if "ego_speed" in df.columns:
                sp_series = pd.to_numeric(df["ego_speed"], errors="coerce")
                if sp_series.notna().any():
                    sp_min, sp_max = float(sp_series.min()), float(sp_series.max())
                    speed_range = st.slider("Ego speed (primary)", sp_min, sp_max, (sp_min, sp_max), step=0.25)
                else:
                    speed_range = None
            else:
                speed_range = None

        with c2:
            # Overlap (secondary) filter if present
            if "overlap" in df.columns:
                ov_series = pd.to_numeric(df["overlap"], errors="coerce").dropna()
                overlaps = sorted([float(x) for x in ov_series.unique()]) if not ov_series.empty else []
                overlaps_pick = st.multiselect("Overlap (secondary)", overlaps, default=overlaps)
            else:
                overlaps_pick = []

        with c3:
            only_pass  = st.checkbox("Show pass only", value=False)
            only_fail  = st.checkbox("Show fail only", value=False)
            only_error = st.checkbox("Show error only", value=False)

    # ---------- Settings (⚙️) ----------
    # Determine candidate numeric columns for tiles
    # Exclude obviously meta/status columns
    exclude_cols = {
        "scenario_name", "status", "state", "pass_flag",
        "date", "aeb_activation", "result_code"
    }
    numeric_candidates = []
    for c in df.columns:
        if c in ("scenario_name", "status", "state", "label1_name", "label2_name", "date"):
            continue
        df[c] = pd.to_numeric(df[c], errors="coerce")
        s = pd.to_numeric(df[c], errors="coerce")
        if s.notna().any():
            numeric_candidates.append(c)

    # Default tile parameters:
    default_tile_params = [x for x in signals_cfg.get("tile_parameters", []) if x in numeric_candidates]
    if not default_tile_params:
        # Reasonable fallback:
        default_tile_params = [x for x in ["ego_speed", "overlap", "impact_speed", "t_aeb"] if x in numeric_candidates]

    with st.expander("Settings ⚙️", expanded=False):
        cA, cB = st.columns(2)
        with cA:
            tile_cols = st.slider("Tile columns", 4, 20, 7, step=1)
        with cB:
            tile_size = st.slider("Tile size", 60, 120, 82, step=2)
        sel_font_px = st.slider("Selected info font size", 10, 32, 20, step=1)

        st.markdown("**Tile parameters (top-to-bottom lines)**")
        tile_fields = st.multiselect(
            "Choose parameters to show on each tile (first one is bold).",
            options=numeric_candidates,
            default=default_tile_params,
            help="Pick any numeric columns from the report. You can add/remove freely."
        )
        st.caption("Units are taken from Signals.json if provided (e.g., `ego_speed`: ` km/h`, `overlap`: ` %`, `impact_speed`: ` m/s`).")

    # Build the mask robustly (fixes your `.all()` crash)
    mask = pd.Series(True, index=df.index)

    if "ego_speed" in df.columns and isinstance(speed_range, (tuple, list, np.ndarray)) and len(speed_range) == 2:
        mask &= pd.to_numeric(df["ego_speed"], errors="coerce").between(speed_range[0], speed_range[1])

    if "overlap" in df.columns and isinstance(overlaps_pick, (list, tuple)) and len(overlaps_pick) > 0:
        mask &= pd.to_numeric(df["overlap"], errors="coerce").isin(overlaps_pick)

    selected_states = []
    if only_pass:  selected_states.append("Pass")
    if only_fail:  selected_states.append("Fail")
    if only_error: selected_states.append("Error")
    if selected_states:
        if "state" in df.columns:
            mask &= df["state"].isin(selected_states)
        elif "status" in df.columns:
            mask &= df["status"].apply(scenario_state_from_status).isin(selected_states)

    view = df[mask].copy()

    # KPIs
    c1, c2, c3, c4 = st.columns(4)
    total = len(view)
    vs = view.copy()
    if "state" not in vs.columns and "status" in vs.columns:
        vs["state"] = vs["status"].apply(scenario_state_from_status)
    passed = int(vs["state"].eq("Pass").sum())
    failed = int(vs["state"].eq("Fail").sum())
    with c1: st.metric("Total Scenarios", total)
    with c2: st.metric("Successful", passed)
    with c3: st.metric("Failed", failed)
    with c4: st.metric("Pass Rate", f"{(passed/total if total else 0):.0%}")

    # Donut + Scenario list (2 columns; list respects global filters)
    col_donut, col_list = st.columns([2, 2])
    with col_donut:
        pie_fig = make_donut(passed, failed)
        st.plotly_chart(pie_fig, use_container_width=True)
    with col_list:
        st.markdown("**Scenario List**")
        render_scenario_list(vs)

    # ---- Tile Panel
    st.subheader("Interactive Scenario Tiles – click a square")
    units_map = signals_cfg.get("units", {}) if isinstance(signals_cfg.get("units"), dict) else {}
    panel = make_tile_panel(view, tile_fields=tile_fields, units=units_map, cols=tile_cols, marker_size=tile_size, label_size=14)
    fig_panel = None
    ordered_df = None

    if panel is None:
        st.info("No scenarios match current filters or primary parameter isn’t present.")
        clicked = []
    else:
        fig_panel, ordered_df = panel

        if plotly_events is None:
            st.warning("Install 'streamlit-plotly-events' to enable tile clicks: pip install streamlit-plotly-events")
            st.plotly_chart(fig_panel, use_container_width=True)
            clicked = []
        else:
            clicked = plotly_events(
                fig_panel,
                click_event=True,
                select_event=False,
                override_height=fig_panel.layout.height or 520,
                override_width="100%"
            )

    # ---------- Selection details (header + 3 columns; all honor font-size) ----------
    details_ph = st.container()  # fixed placeholder so table below never disappears

    with details_ph:
        if clicked and ordered_df is not None:
            idx = clicked[0].get("pointIndex", clicked[0].get("pointNumber"))
            if idx is not None and 0 <= idx < len(ordered_df):
                sel_row = ordered_df.iloc[int(idx)]

                # 1) Header "Selected:"
                header_html = (
                    f'<div style="font-size:{sel_font_px}px; line-height:1.3">'
                    f'<b>Selected:</b> <code>{sel_row.get("scenario_name","")}</code>'
                )
                if "ego_speed" in sel_row and pd.notna(sel_row["ego_speed"]):
                    header_html += f' &nbsp; | &nbsp; {sel_row.get("ego_speed","")} kmph'
                if "overlap" in sel_row and pd.notna(sel_row["overlap"]):
                    header_html += f', {sel_row.get("overlap","")}%'
                header_html += f' &nbsp; | &nbsp; Status: <b>{sel_row.get("status","")}</b></div>'
                st.markdown(header_html, unsafe_allow_html=True)

                # 2) 3-column detail layout (guard missing columns)
                c1, c2, c3 = st.columns(3)
                with c1:
                    st.markdown(
                        f'<div style="font-size:{sel_font_px}px; line-height:1.5">'
                        f'<b>Impact speed:</b> {sel_row.get("impact_speed","")} m/s<br>'
                        f'<b>Stop distance:</b> {sel_row.get("stop_distance","")} m'
                        f'</div>', unsafe_allow_html=True
                    )
                with c2:
                    st.markdown(
                        f'<div style="font-size:{sel_font_px}px; line-height:1.5">'
                        f'<b>T_AEB:</b> {sel_row.get("t_aeb","")} s<br>'
                        f'<b>AEB activation:</b> {sel_row.get("aeb_activation","")}'
                        f'</div>', unsafe_allow_html=True
                    )
                with c3:
                    st.markdown(
                        f'<div style="font-size:{sel_font_px}px; line-height:1.5">'
                        f'<b>Result code:</b> {sel_row.get("result_code","")}<br>'
                        f'<b>Date:</b> {sel_row.get("date","")}'
                        f'</div>', unsafe_allow_html=True
                    )

    st.divider()

    # ---- All Scenarios Table (robust + stays visible) ----
    st.subheader("All Scenarios")
    all_ph = st.container()   # fixed placeholder so it persists after clicks

    with all_ph:
        show_cols = [c for c in [
            "scenario_name", "state", "status", "ego_speed", "overlap",
            "impact_speed", "stop_distance", "t_aeb", "aeb_activation",
            "result_code", "date"
        ] if c in view.columns]
        table = view[show_cols]

        if table.empty:
            st.info("No rows to display for the current filters.")
        else:
            if USE_AGGRID:
                try:
                    g2 = GridOptionsBuilder.from_dataframe(table)
                    g2.configure_default_column(filter=True, sortable=True, resizable=True)
                    g2.configure_selection("single")
                    AgGrid(
                        table,
                        gridOptions=g2.build(),
                        update_mode=GridUpdateMode.NO_UPDATE,
                        height=360,
                        fit_columns_on_grid_load=True,
                        key="all_scenarios_grid"   # stable key
                    )
                except Exception as ex:
                    st.warning(f"AgGrid error ({ex}). Falling back to basic table.")
                    USE_AGGRID_LOCAL = False
                else:
                    USE_AGGRID_LOCAL = True
            else:
                USE_AGGRID_LOCAL = False

            if not USE_AGGRID_LOCAL:
            
                # ---- CLOUD-SAFE fallback (NO pandas Styler) ----
                table_view = table.copy()
            
                # Drop empty columns/rows (Arrow-safe)
                table_view = table_view.dropna(axis=1, how="all")
                table_view = table_view.dropna(axis=0, how="all")
            
                # Remove duplicate columns (CRITICAL for Streamlit Cloud)
                table_view = table_view.loc[:, ~table_view.columns.duplicated()]
            
                st.dataframe(
                    table_view,
                    use_container_width=True,
                    hide_index=True
                )
            

    # ---- Export PPT ----
    st.subheader("Export")
    if st.button("Export PowerPoint (.pptx)"):
        if total == 0:
            st.warning("Nothing to export for the current filters.")
        else:
            out = io.BytesIO()
            export_ppt(view_df=view, pie_fig=make_donut(passed, failed), panel_fig=fig_panel, file_bytes=out)

            # 1) Download (works in browser mode)
            st.success("Report generated. You can download it below or find it in the 'exports' folder.")
            st.download_button(
                "Download PPTX",
                data=out.getvalue(),
                file_name="ADAS_Batch_Report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

            # 2) Also save to disk (useful when running inside pywebview desktop)
            ts = pd.Timestamp.now().strftime("%Y%m%d_%H%M%S")
            export_dir = os.path.join(os.getcwd(), "exports")
            os.makedirs(export_dir, exist_ok=True)
            file_path = os.path.join(export_dir, f"ADAS_Batch_Report_{ts}.pptx")
            with open(file_path, "wb") as f:
                f.write(out.getvalue())
            st.info(f"Saved to: `{file_path}`")

    # ---- Footer: Help + About ----
    st.divider()
    ##Help: read Batch_readme.md if present
    # with st.expander("Help • Readme (Batch_readme.md)", expanded=False):
    #     readme_path = os.path.join(os.getcwd(), "Batch_readme.md")
    #     if os.path.exists(readme_path):
    #         try:
    #             with open(readme_path, "r", encoding="utf-8", errors="ignore") as f:
    #                 st.markdown(f.read())
    #         except Exception as ex:
    #             st.warning(f"Could not read Batch_readme.md ({ex})")
    #     else:
    #         st.info("Place a file named `Batch_readme.md` next to this app to show contextual help.")

    # About: small popup if available, fallback to expander
    try:
        pop = st.popover("About")
        with pop:
            st.markdown(
                "**ADAS Batch Viewer**  \n"
                "Version: 1.1  \n"
                "Last update on 10th March 2026.  \n"
                "Created by Amal Nambiar for Mahindra & Mahindra.\n  \n"
                "© 2026"
            )
    except Exception:
        with st.expander("About", expanded=False):
            st.markdown(
                "**ADAS Batch Viewer**  \n"
                "Version: 1.0  \n"
                "Last update on 10th March 2026.  \n"
                "Created by Amal Nambiar for Mahindra & Mahindra.\n  \n"
                "© 2026"
            )

else:
    st.info("Upload your batch `index.html` to begin.")
